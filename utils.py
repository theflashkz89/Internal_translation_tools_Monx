import os
import json
import streamlit as st
from openai import OpenAI
from typing import Tuple, Optional, List, Dict, Any
import PyPDF2
from pdf2docx import Converter
from docx import Document
import uuid
from pathlib import Path
import time
import re
from concurrent.futures import ThreadPoolExecutor, as_completed
from threading import Lock
import deepl

# 全局 DeepL 翻译器实例（只初始化一次）
deepl_translator = None

def _get_deepl_translator():
    """
    获取全局 DeepL 翻译器实例（单例模式）
    """
    global deepl_translator
    
    if deepl_translator is None:
        # 显式初始化变量，防止报错
        api_key = None
        
        # 尝试从 Streamlit secrets 获取
        try:
            if hasattr(st, 'secrets') and hasattr(st.secrets, 'get'):
                api_key = st.secrets.get("DEEPL_API_KEY")
        except Exception:
            pass
        
        # 如果 secrets 中没有，尝试从环境变量获取
        if not api_key:
            api_key = os.getenv("DEEPL_API_KEY")
        
        if not api_key:
            raise ValueError(
                "未找到 DeepL API Key。请配置环境变量 DEEPL_API_KEY 或 "
                "在 Streamlit secrets 中设置 DEEPL_API_KEY"
            )
        
        # 根据 Key 后缀判断使用哪个 URL
        if api_key.endswith(':fx'):
            deepl_translator = deepl.Translator(
                auth_key=api_key,
                server_url="https://api-free.deepl.com"
            )
        else:
            deepl_translator = deepl.Translator(auth_key=api_key)
    
    return deepl_translator

# 术语库：从 glossary.json 加载
# 用于减少 API 调用并提高翻译准确度
def load_glossary() -> dict:
    """
    从根目录下的 glossary.json 文件加载术语库
    
    返回:
        dict: 术语库字典（中文 -> 英文），如果文件不存在或格式错误则返回空字典
    """
    glossary_path = Path("glossary.json")
    
    try:
        if not glossary_path.exists():
            print(f"警告：术语库文件 {glossary_path} 不存在，使用默认术语库")
            return {}
        
        with open(glossary_path, 'r', encoding='utf-8') as f:
            glossary = json.load(f)
            
        if not isinstance(glossary, dict):
            print(f"警告：术语库文件 {glossary_path} 格式错误（不是有效的JSON对象），使用默认术语库")
            return {}
        
        return glossary
        
    except json.JSONDecodeError as e:
        print(f"警告：术语库文件 {glossary_path} JSON格式错误：{str(e)}，使用默认术语库")
        return {}
    except Exception as e:
        print(f"警告：加载术语库文件 {glossary_path} 时发生错误：{str(e)}，使用默认术语库")
        return {}


# 在模块加载时初始化术语库
GLOSSARY = load_glossary()


def call_deepseek_api(text: str, prompt: str, model: str = "deepseek-chat", max_retries: int = 3) -> str:
    """
    调用 DeepSeek API 获取 AI 回复（带重试机制）
    
    参数:
        text (str): 用户输入的文本
        prompt (str): 系统提示词/指令
        model (str): 使用的模型名称，默认为 "deepseek-chat"
        max_retries (int): 最大重试次数，默认为3次
    
    返回:
        str: AI 的回复内容
    
    异常:
        Exception: 当 API 调用失败时抛出异常
    """
    # 从 Streamlit secrets 或环境变量获取 API Key
    api_key = None
    try:
        # 尝试从 Streamlit secrets 获取
        if hasattr(st, 'secrets') and hasattr(st.secrets, 'get'):
            api_key = st.secrets.get("DEEPSEEK_API_KEY")
    except:
        pass
    
    # 如果 secrets 中没有，尝试从环境变量获取
    if not api_key:
        api_key = os.getenv("DEEPSEEK_API_KEY")
    
    if not api_key:
        raise ValueError(
            "未找到 DeepSeek API Key。请配置环境变量 DEEPSEEK_API_KEY 或 "
            "在 Streamlit secrets 中设置 DEEPSEEK_API_KEY"
        )
    
    # 初始化 OpenAI 客户端，使用 DeepSeek 的 base_url
    client = OpenAI(
        api_key=api_key,
        base_url="https://api.deepseek.com"
    )
    
    # 重试机制
    last_exception = None
    for attempt in range(max_retries):
        try:
            # 调用 API，设置超时时间
            response = client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": prompt},
                    {"role": "user", "content": text}
                ],
                temperature=0.7,
                timeout=60.0,  # 设置60秒超时
            )
            
            # 返回 AI 的回复
            return response.choices[0].message.content
        
        except Exception as e:
            last_exception = e
            error_msg = str(e)
            
            # 如果是认证错误，不需要重试
            if "401" in error_msg or "unauthorized" in error_msg.lower():
                raise Exception(f"认证失败：API Key 无效或已过期。请检查 .streamlit/secrets.toml 中的配置。")
            
            # 如果是速率限制，等待更长时间后重试
            if "429" in error_msg or "rate limit" in error_msg.lower():
                if attempt < max_retries - 1:
                    wait_time = (attempt + 1) * 5  # 递增等待时间：5秒、10秒、15秒
                    time.sleep(wait_time)
                    continue
                else:
                    raise Exception(f"请求过于频繁：已重试{max_retries}次，请稍后再试。")
            
            # 如果是超时或连接错误，等待后重试
            if ("timeout" in error_msg.lower() or "timed out" in error_msg.lower() or 
                "connection" in error_msg.lower() or "connect" in error_msg.lower()):
                if attempt < max_retries - 1:
                    wait_time = (attempt + 1) * 2  # 递增等待时间：2秒、4秒、6秒
                    time.sleep(wait_time)
                    continue
                else:
                    if "timeout" in error_msg.lower() or "timed out" in error_msg.lower():
                        raise Exception(f"连接超时：API 请求超过60秒未响应。已重试{max_retries}次，请检查网络连接后重试。")
                    else:
                        raise Exception(f"连接失败：无法连接到 DeepSeek API 服务器。已重试{max_retries}次。\n详细错误：{error_msg}")
            
            # 其他错误，如果是最后一次尝试，抛出异常
            if attempt == max_retries - 1:
                raise Exception(f"调用 DeepSeek API 时发生错误：{error_msg}")
            else:
                # 其他错误也等待后重试
                time.sleep(2)
                continue
    
    # 如果所有重试都失败
    raise Exception(f"调用 DeepSeek API 失败：{str(last_exception)}")


def generate_email_draft(
    email_type: str,
    tone: str,
    language: str,
    recipient: str,
    subject: str,
    key_points: str
) -> str:
    """
    使用 DeepSeek API 生成邮件草稿
    
    参数:
        email_type (str): 邮件类型（商务邮件、感谢信、请求邮件、通知邮件、回复邮件）
        tone (str): 语气风格（正式、友好、简洁、礼貌）
        language (str): 邮件语言（中文、英文、意大利语）
        recipient (str): 收件人称呼（可选）
        subject (str): 邮件主题
        key_points (str): 关键要点/背景信息
    
    返回:
        str: 生成的邮件草稿正文
    
    注意:
        当 language="意大利语" 且 tone="正式" 时，会自动使用意大利语正式尊称（Lei形式）
    """
    # 构建系统提示词
    # 检查是否需要使用意大利语尊称（正式语气 + 意大利语）
    italian_formal_note = ""
    if language == "意大利语" and tone == "正式":
        italian_formal_note = "- 必须使用意大利语正式尊称：使用 Lei（您，大写L）代替 tu，所有动词使用第三人称单数形式（如：Lei è, Lei può, Lei desidera 等），这是意大利语正式邮件的标准礼仪\n"
    
    system_prompt = f"""你是一位专业的商务写作助手。
请根据用户提供的信息，撰写一封{email_type}。
要求：
- 语气：{tone}
- 语言：{language}
- 格式规范，结构清晰
- 只输出邮件正文，不要添加额外说明
- 如果是中文邮件，使用适当的敬语和礼貌用语
- 如果是英文邮件，遵循商务邮件的标准格式
{italian_formal_note}- 重要：不要使用任何 Markdown 格式符号，包括星号（*）、下划线（_）、井号（#）等
- 输出纯文本格式的邮件正文，不要有任何格式标记符号"""
    
    # 构建用户输入
    recipient_text = recipient if recipient else "（未指定）"
    user_input = f"""请撰写一封邮件：
收件人称呼：{recipient_text}
邮件主题：{subject}
关键要点：
{key_points}"""
    
    return call_deepseek_api(user_input, system_prompt)


def _should_translate_text(text: str) -> bool:
    """
    检查文本是否需要翻译
    
    参数:
        text (str): 要检查的文本
    
    返回:
        bool: True表示需要翻译，False表示可以跳过
    """
    text_stripped = text.strip()
    
    # 如果为空，不需要翻译
    if not text_stripped:
        return False
    
    # 检查是否为纯数字（包括小数、负数、科学计数法等）
    # 匹配：纯数字、带小数点的数字、负数、百分比、货币等
    if re.match(r'^[\d\s\.,\-+\%\$€¥]+$', text_stripped):
        # 进一步检查是否包含至少一个数字
        if re.search(r'\d', text_stripped):
            return False
    
    # 检查是否为日期格式（如：2025-12-24, 2025/12/24, 2025.12.24等）
    date_patterns = [
        r'^\d{4}[-/\.]\d{1,2}[-/\.]\d{1,2}$',  # 2025-12-24, 2025/12/24
        r'^\d{1,2}[-/\.]\d{1,2}[-/\.]\d{4}$',  # 12-24-2025
        r'^\d{4}年\d{1,2}月\d{1,2}日$',  # 2025年12月24日
    ]
    for pattern in date_patterns:
        if re.match(pattern, text_stripped):
            return False
    
    # 检查是否只包含标点符号和空白字符
    if re.match(r'^[\s\.,;:!?\-_()\[\]{}"\']+$', text_stripped):
        return False
    
    # 其他情况需要翻译
    return True


def translate_text(text: str, target_language: str) -> str:
    """
    翻译文本到目标语言
    
    参数:
        text (str): 要翻译的文本
        target_language (str): 目标语言名称（如：英语、日语、法语等）
    
    返回:
        str: 翻译后的文本
    
    异常:
        Exception: 当翻译失败时抛出异常
    """
    # 获取目标语言代码
    target_lang_code = _get_deepl_lang_code(target_language)
    
    # 检查术语库：只有当目标语言是英语时，才使用术语库
    # 术语库仅支持中文→英文方向
    text_stripped = text.strip()
    if target_lang_code.upper().startswith("EN") and text_stripped in GLOSSARY:
        return GLOSSARY[text_stripped]
    
    # 如果不在术语库中或目标语言不是英语，使用 DeepL API 进行翻译
    return call_deepl_api(text, target_lang=target_lang_code)


def _get_deepl_lang_code(language_name: str) -> str:
    """
    将语言名称转换为 DeepL 语言代码
    
    参数:
        language_name (str): 语言名称（如：英语、日语、法语等）
    
    返回:
        str: DeepL 语言代码（如：EN-US、JA、FR 等）
    """
    language_map = {
        "中文": "ZH",
        "简体中文": "ZH",
        "英语": "EN-US",
        "英文": "EN-US",  # 添加"英文"映射
        "日语": "JA",
        "法语": "FR",
        "德语": "DE",
        "西班牙语": "ES",
        "俄语": "RU",
        "韩语": "KO",
        "意大利语": "IT",
        "葡萄牙语": "PT",
        "阿拉伯语": "AR",
        "泰语": "TH",
        "越南语": "VI",
        "印尼语": "ID",
        "荷兰语": "NL",
        "瑞典语": "SV",
        "挪威语": "NO",
        "丹麦语": "DA",
        "芬兰语": "FI",
        "波兰语": "PL",
        "土耳其语": "TR"
    }
    # 默认返回英语（美式）
    return language_map.get(language_name, "EN-US")


def call_deepl_api(text: str, target_lang: str = "EN-US") -> str:
    """
    使用 DeepL API 翻译文本（单条文本）
    
    参数:
        text (str): 要翻译的文本
        target_lang (str): 目标语言代码，默认为 "EN-US"（美式英语）
            - EN-US: 美式英语
            - EN-GB: 英式英语
            - 其他 DeepL 支持的语言代码
    
    返回:
        str: 翻译后的文本
    
    异常:
        ValueError: 当 API Key 未配置时抛出异常
        Exception: 当 DeepL API 调用失败时抛出异常（包括额度用完、网络错误等）
    
    说明:
        - 优先检查本地术语库（GLOSSARY），如果存在则直接返回，不消耗 DeepL 额度
        - 术语库仅适用于翻译到英语（target_lang 以 "EN" 开头）
        - 如果不在术语库中或目标语言不是英语，则调用 DeepL API 进行翻译
    """
    # 检查术语库：只有当目标语言是英语时，才使用术语库
    text_stripped = text.strip()
    if target_lang.upper().startswith("EN") and text_stripped in GLOSSARY:
        return GLOSSARY[text_stripped]
    
    # 如果不在术语库中，调用 DeepL API
    try:
        # 获取 DeepL 翻译器实例
        translator = _get_deepl_translator()
        
        # 调用 DeepL API 进行翻译
        result = translator.translate_text(text, target_lang=target_lang)
        
        # 返回翻译结果
        return result.text
        
    except ValueError as e:
        # API Key 未配置的错误已经在 _get_deepl_translator() 中处理
        raise e
    except deepl.exceptions.QuotaExceededException:
        raise Exception("DeepL API 额度已用完。请检查您的账户余额或升级套餐。")
    except deepl.exceptions.AuthorizationException:
        raise Exception("DeepL API 认证失败。请检查 API Key 是否正确。")
    except deepl.exceptions.TooManyRequestsException:
        raise Exception("DeepL API 请求过于频繁。请稍后再试。")
    except deepl.exceptions.ConnectionException as e:
        raise Exception(f"DeepL API 连接失败：{str(e)}。请检查网络连接。")
    except deepl.exceptions.DeepLException as e:
        raise Exception(f"DeepL API 调用失败：{str(e)}")
    except Exception as e:
        # 捕获其他未知异常
        raise Exception(f"调用 DeepL API 时发生未知错误：{str(e)}")


def call_deepl_api_batch(texts: List[str], target_lang: str = "EN-US") -> List[str]:
    """
    使用 DeepL API 批量翻译文本列表（批处理模式，更高效）
    
    参数:
        texts (list[str]): 要翻译的文本列表
        target_lang (str): 目标语言代码，默认为 "EN-US"（美式英语）
    
    返回:
        list[str]: 翻译后的文本列表，顺序与输入列表一致
    
    异常:
        ValueError: 当 API Key 未配置时抛出异常
        Exception: 当 DeepL API 调用失败时抛出异常（包括额度用完、网络错误等）
    
    说明:
        - 优先检查本地术语库（GLOSSARY），如果存在则直接返回，不消耗 DeepL 额度
        - 术语库仅适用于翻译到英语（target_lang 以 "EN" 开头）
        - 如果不在术语库中或目标语言不是英语，则调用 DeepL API 进行批量翻译
        - 批处理模式比单条翻译更高效，减少网络开销
    """
    if not texts:
        return []
    
    # 检查目标语言是否为英语（术语库仅适用于中译英）
    is_target_english = target_lang.upper().startswith("EN")
    
    # 分离需要从术语库获取的文本和需要 API 翻译的文本
    glossary_results = {}  # {索引: 翻译文本}
    api_texts = []  # 需要 API 翻译的文本列表
    api_indices = []  # 对应的原始索引列表
    
    for idx, text in enumerate(texts):
        text_stripped = text.strip()
        # 只有当目标语言是英语时，才使用术语库
        if is_target_english and text_stripped in GLOSSARY:
            glossary_results[idx] = GLOSSARY[text_stripped]
        else:
            api_texts.append(text)
            api_indices.append(idx)
    
    # 如果所有文本都在术语库中，直接返回
    if not api_texts:
        return [glossary_results.get(i, texts[i]) for i in range(len(texts))]
    
    # 调用 DeepL API 批量翻译
    try:
        translator = _get_deepl_translator()
        
        # 批量翻译
        results = translator.translate_text(api_texts, target_lang=target_lang)
        
        # 处理结果：results 可能是单个结果或结果列表
        if isinstance(results, list):
            api_translations = [r.text for r in results]
        else:
            # 如果只有一个结果，DeepL 可能返回单个对象
            api_translations = [results.text]
        
        # 构建完整的翻译结果列表
        final_results = []
        api_idx = 0
        for i in range(len(texts)):
            if i in glossary_results:
                final_results.append(glossary_results[i])
            else:
                if api_idx < len(api_translations):
                    final_results.append(api_translations[api_idx])
                    api_idx += 1
                else:
                    # 如果结果数量不匹配，保留原文
                    final_results.append(texts[i])
        
        return final_results
        
    except ValueError as e:
        raise e
    except deepl.exceptions.QuotaExceededException:
        raise Exception("DeepL API 额度已用完。请检查您的账户余额或升级套餐。")
    except deepl.exceptions.AuthorizationException:
        raise Exception("DeepL API 认证失败。请检查 API Key 是否正确。")
    except deepl.exceptions.TooManyRequestsException:
        raise Exception("DeepL API 请求过于频繁。请稍后再试。")
    except deepl.exceptions.ConnectionException as e:
        raise Exception(f"DeepL API 连接失败：{str(e)}。请检查网络连接。")
    except deepl.exceptions.DeepLException as e:
        raise Exception(f"DeepL API 调用失败：{str(e)}")
    except Exception as e:
        raise Exception(f"调用 DeepL API 时发生未知错误：{str(e)}")


def handle_pdf_processing(pdf_file) -> Tuple[Optional[str], Optional[str]]:
    """
    处理PDF文件：检查是否可提取文本，并转换为Word格式
    
    参数:
        pdf_file: 用户上传的PDF文件对象（Streamlit UploadedFile对象）
    
    返回:
        Tuple[Optional[str], Optional[str]]: 
            - 如果成功：返回 (生成的Word文件路径, None)
            - 如果失败：返回 (None, 错误信息)
    
    异常处理:
        - 如果PDF是扫描版（无法提取文本），返回错误信息
        - 如果转换失败，抛出异常并返回错误信息
    """
    try:
        # 步骤1：检查PDF是否包含可识别文本
        # 将文件指针重置到开头
        pdf_file.seek(0)
        
        # 使用PyPDF2读取PDF文件
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        # 检查PDF是否有页面
        if len(pdf_reader.pages) == 0:
            return None, "PDF文件为空，无法处理"
        
        # 读取第一页并提取文本
        first_page = pdf_reader.pages[0]
        extracted_text = first_page.extract_text()
        
        # 检查提取的文本长度（少于10个字符认为是扫描版）
        if len(extracted_text.strip()) < 10:
            return None, "检测到是扫描版PDF，暂不支持"
        
        # 步骤2：转换PDF为Word格式
        # 创建temp文件夹（如果不存在）
        temp_dir = Path("temp")
        temp_dir.mkdir(exist_ok=True)
        
        # 生成唯一的临时文件名
        unique_id = str(uuid.uuid4())
        temp_pdf_path = temp_dir / f"temp_{unique_id}.pdf"
        output_docx_path = temp_dir / f"converted_{unique_id}.docx"
        
        # 将上传的文件保存到临时位置
        pdf_file.seek(0)  # 重置文件指针
        with open(temp_pdf_path, "wb") as f:
            f.write(pdf_file.read())
        
        # 使用pdf2docx进行转换
        try:
            cv = Converter(str(temp_pdf_path))
            cv.convert(str(output_docx_path))
            cv.close()
        except Exception as e:
            # 清理临时文件
            if temp_pdf_path.exists():
                temp_pdf_path.unlink()
            raise Exception(f"PDF转Word失败：{str(e)}")
        finally:
            # 删除临时PDF文件
            if temp_pdf_path.exists():
                temp_pdf_path.unlink()
        
        # 返回生成的Word文件路径
        return str(output_docx_path), None
        
    except Exception as e:
        error_msg = str(e)
        # 如果是我们自定义的错误信息，直接返回
        if "检测到是扫描版PDF" in error_msg or "PDF文件为空" in error_msg:
            return None, error_msg
        # 其他异常，返回通用错误信息
        return None, f"处理PDF文件时发生错误：{error_msg}"


def translate_word_document(docx_path: str, target_language: str, progress_callback=None) -> str:
    """
    翻译Word文档中的所有段落文本和表格内容（使用多线程并发模式优化性能）
    
    参数:
        docx_path (str): Word文件的路径
        target_language (str): 目标语言名称（如：英语、日语、法语等）
        progress_callback (callable, optional): 进度回调函数，接收 (current, total, status) 参数
    
    返回:
        str: 翻译后保存的新文件路径
    
    异常:
        Exception: 当文件读取、翻译或保存失败时抛出异常
    
    处理逻辑:
        1. 使用python-docx加载Word文档
        2. 收集所有需要翻译的段落（正文+表格）
        3. 先处理术语库中的文本（直接应用，不占用线程池）
        4. 对需要API翻译的文本，使用ThreadPoolExecutor并发执行（max_workers=5）
        5. 确保翻译结果按顺序填回文档
        6. 将翻译好的文档保存为新文件
    """
    # 检查目标语言是否为英语（术语库仅适用于中译英）
    is_target_english = "英" in target_language or "English" in target_language
    
    # 术语库已在模块加载时初始化，无需重复加载
    
    def _apply_translation_to_paragraph(paragraph, translated_text: str):
        """
        将翻译后的文本应用到段落对象
        
        参数:
            paragraph: python-docx 的段落对象
            translated_text (str): 翻译后的文本
        """
        paragraph.clear()
        paragraph.add_run(translated_text)
    
    def _translate_batch_task(batch_data: Tuple[List, str]) -> Tuple[List, List[str], Optional[Exception]]:
        """
        批量翻译任务的辅助函数（用于线程池）
        
        参数:
            batch_data: (任务数据列表, 目标语言代码) 的元组
                任务数据列表格式：[(任务索引, 待翻译文本, paragraph对象), ...]
        
        返回:
            Tuple[list, list[str], Optional[Exception]]: 
                (任务数据列表, 翻译结果列表, 异常对象)
                如果成功：返回 (任务数据列表, 翻译文本列表, None)
                如果失败：返回 (任务数据列表, [], 异常对象)
        """
        task_list, target_lang_code = batch_data
        
        try:
            # 提取待翻译的文本列表
            texts = [task[1] for task in task_list]
            
            # 调用 DeepL 批量翻译 API（不需要 Prompt）
            translated_texts = call_deepl_api_batch(texts, target_lang_code)
            
            return (task_list, translated_texts, None)
        except Exception as e:
            # 捕获异常，返回异常对象
            return (task_list, [], e)
    
    try:
        # 加载Word文档
        if progress_callback:
            progress_callback(0, 0, "正在加载文档...")
        
        doc = Document(docx_path)
        
        # 分析文档结构
        if progress_callback:
            progress_callback(0, 0, "正在分析文档结构...")
        
        # ========== 第一步：收集所有待翻译任务 ==========
        # 任务格式：(任务索引, 待翻译文本, paragraph对象, 目标语言)
        translation_tasks = []
        
        # 收集正文段落
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text and len(text) <= 8000:  # 跳过超长段落
                # 检查是否需要翻译（跳过纯数字、日期等）
                if _should_translate_text(text):
                    translation_tasks.append((len(translation_tasks), text, paragraph, target_language))
        
        # 收集表格单元格中的段落
        for table_idx, table in enumerate(doc.tables):
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    for para_idx, paragraph in enumerate(cell.paragraphs):
                        text = paragraph.text.strip()
                        if text and len(text) <= 8000:  # 跳过超长段落
                            # 检查是否需要翻译（跳过纯数字、日期等）
                            if _should_translate_text(text):
                                translation_tasks.append((len(translation_tasks), text, paragraph, target_language))
                            else:
                                # 不需要翻译的段落（纯数字、日期等），保留原文
                                pass
        
        total_count = len(translation_tasks)
        processed = 0
        failed = 0
        
        if progress_callback:
            progress_callback(0, total_count, f"开始翻译，共 {total_count} 个段落...")
        
        # ========== 第二步：先处理术语库中的文本（直接应用，不占用线程池）==========
        api_tasks = []  # 需要API翻译的任务列表
        glossary_results = {}  # 术语库翻译结果：{任务索引: 翻译文本}
        
        for task_idx, text, paragraph, _ in translation_tasks:
            text_stripped = text.strip()
            # 只有当目标语言是英语时，才使用术语库
            if is_target_english and text_stripped in GLOSSARY:
                # 命中术语库：直接应用翻译
                _apply_translation_to_paragraph(paragraph, GLOSSARY[text_stripped])
                glossary_results[task_idx] = GLOSSARY[text_stripped]
                processed += 1
            else:
                # 未命中术语库或目标语言不是英语：加入API翻译任务列表
                api_tasks.append((task_idx, text, paragraph, target_language))
        
        # ========== 第三步：将任务分批并使用线程池并发执行API翻译 ==========
        if api_tasks:
            # 将语言名称转换为 DeepL 语言代码
            target_lang_code = _get_deepl_lang_code(target_language)
            
            # 将任务分批（每 50 个一组）
            batch_size = 50
            task_batches = []
            for i in range(0, len(api_tasks), batch_size):
                batch = api_tasks[i:i + batch_size]
                task_batches.append((batch, target_lang_code))
            
            if progress_callback:
                progress_callback(processed, total_count, f"正在并发翻译 {len(api_tasks)} 个段落（分为 {len(task_batches)} 批，使用 {min(5, len(task_batches))} 个线程）...")
            
            # 用于线程安全的进度更新
            progress_lock = Lock()
            completed_batches = [0]  # 使用列表以便在闭包中修改
            
            # 存储翻译结果：{任务索引: (翻译文本, 异常)}
            translation_results = {}
            
            # 使用ThreadPoolExecutor并发执行批处理任务
            max_workers = 5  # 最大并发数
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                # 提交所有批处理任务
                future_to_batch = {
                    executor.submit(_translate_batch_task, batch_data): batch_data 
                    for batch_data in task_batches
                }
                
                # 使用as_completed来获取完成的任务（不保证顺序）
                for future in as_completed(future_to_batch):
                    batch_data = future_to_batch[future]
                    task_list = batch_data[0]
                    
                    try:
                        result_task_list, translated_texts, error = future.result()
                        
                        if error is None:
                            # 批处理翻译成功
                            for idx, (task_idx, _, paragraph, _) in enumerate(result_task_list):
                                if idx < len(translated_texts):
                                    translation_results[task_idx] = (translated_texts[idx], None)
                                else:
                                    # 结果数量不匹配，标记为失败
                                    translation_results[task_idx] = (None, Exception("批处理结果数量不匹配"))
                        else:
                            # 批处理翻译失败，标记该批次所有任务为失败
                            for task_idx, _, _, _ in task_list:
                                translation_results[task_idx] = (None, error)
                            
                            # 检查是否是速率限制错误，如果是则等待后重试
                            error_msg = str(error)
                            if "429" in error_msg or "rate limit" in error_msg.lower() or "TooManyRequests" in error_msg:
                                # 等待一段时间后重试（在主线程中重试，避免阻塞线程池）
                                time.sleep(5)
                                try:
                                    retry_task_list, retry_texts, retry_error = _translate_batch_task(batch_data)
                                    if retry_error is None:
                                        # 重试成功，覆盖失败结果
                                        for idx, (task_idx, _, _, _) in enumerate(retry_task_list):
                                            if idx < len(retry_texts):
                                                translation_results[task_idx] = (retry_texts[idx], None)
                                    else:
                                        # 重试也失败，保留失败结果
                                        pass
                                except Exception as retry_ex:
                                    # 重试时发生异常，保留原始错误
                                    pass
                        
                        # 更新进度（线程安全）
                        with progress_lock:
                            completed_batches[0] += 1
                            if progress_callback:
                                current_processed = processed + sum(1 for tid in translation_results if translation_results[tid][1] is None)
                                progress_callback(
                                    current_processed, 
                                    total_count, 
                                    f"翻译进度：{current_processed}/{total_count} ({completed_batches[0]}/{len(task_batches)} 批已完成)"
                                )
                    
                    except Exception as e:
                        # 处理future异常，标记该批次所有任务为失败
                        for task_idx, _, _, _ in task_list:
                            translation_results[task_idx] = (None, e)
                        with progress_lock:
                            completed_batches[0] += 1
                            if progress_callback:
                                current_processed = processed + sum(1 for tid in translation_results if translation_results[tid][1] is None)
                                progress_callback(
                                    current_processed, 
                                    total_count, 
                                    f"翻译进度：{current_processed}/{total_count} (批次失败)"
                                )
            
            # ========== 第四步：按顺序应用翻译结果 ==========
            if progress_callback:
                progress_callback(processed, total_count, "正在应用翻译结果...")
            
            # 按任务索引顺序应用结果
            for task_idx, text, paragraph, _ in api_tasks:
                if task_idx in translation_results:
                    translated_text, error = translation_results[task_idx]
                    if error is None and translated_text:
                        # 翻译成功，应用结果
                        _apply_translation_to_paragraph(paragraph, translated_text)
                        processed += 1
                    else:
                        # 翻译失败，保留原文
                        failed += 1
                else:
                    # 结果缺失，保留原文
                    failed += 1
        
        # 生成输出文件路径
        if progress_callback:
            progress_callback(total_count, total_count, "正在保存翻译后的文档...")
        
        original_path = Path(docx_path)
        output_path = original_path.parent / f"translated_{original_path.stem}.docx"
        
        # 保存翻译后的文档
        doc.save(str(output_path))
        
        if progress_callback:
            if failed > 0:
                progress_callback(total_count, total_count, f"翻译完成！成功：{processed} 个，跳过/失败：{failed} 个")
            else:
                progress_callback(total_count, total_count, f"翻译完成！共翻译 {processed} 个段落")
        
        return str(output_path)
        
    except FileNotFoundError:
        raise Exception(f"找不到文件：{docx_path}")
    except Exception as e:
        raise Exception(f"翻译Word文档时发生错误：{str(e)}")


def apply_custom_styles():
    """
    应用自定义样式（DeepL风格）
    使用多阶段加载策略防止UI闪烁：
    1. 第一阶段：立即隐藏页面，应用关键样式
    2. 第二阶段：加载完整样式
    3. 第三阶段：渐入显示页面
    """
    # 自定义连接断开提示脚本（隐藏默认弹窗，显示友好提示）
    import streamlit.components.v1 as components
    
    connection_monitor_html = """
    <style>
        /* 隐藏 Streamlit 默认的连接错误弹窗 */
        .stException,
        div[data-testid="stConnectionStatus"],
        .element-container:has(div[data-testid="stConnectionStatus"]) {
            display: none !important;
        }
        
        /* 自定义重连提示样式 */
        #custom-reconnect-overlay {
            display: none;
            position: fixed;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            background: rgba(255, 255, 255, 0.95);
            z-index: 999999;
            justify-content: center;
            align-items: center;
            flex-direction: column;
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
        }
        
        #custom-reconnect-overlay.show {
            display: flex !important;
        }
        
        #custom-reconnect-content {
            text-align: center;
            padding: 40px;
            max-width: 400px;
        }
        
        #custom-reconnect-icon {
            font-size: 48px;
            margin-bottom: 20px;
        }
        
        #custom-reconnect-title {
            font-size: 24px;
            font-weight: 600;
            color: #0F2B46;
            margin-bottom: 12px;
        }
        
        #custom-reconnect-message {
            font-size: 16px;
            color: #64748B;
            margin-bottom: 24px;
            line-height: 1.5;
        }
        
        #custom-reconnect-btn {
            background: linear-gradient(135deg, #3B82F6 0%, #2563EB 100%);
            color: white;
            border: none;
            padding: 14px 32px;
            font-size: 16px;
            font-weight: 600;
            border-radius: 8px;
            cursor: pointer;
            transition: all 0.2s ease;
            box-shadow: 0 4px 12px rgba(59, 130, 246, 0.35);
        }
        
        #custom-reconnect-btn:hover {
            transform: translateY(-2px);
            box-shadow: 0 6px 16px rgba(59, 130, 246, 0.45);
        }
        
        #custom-reconnect-btn:active {
            transform: translateY(0);
        }
    </style>
    
    <!-- 自定义重连提示 -->
    <div id="custom-reconnect-overlay">
        <div id="custom-reconnect-content">
            <div id="custom-reconnect-icon">☕</div>
            <div id="custom-reconnect-title">您已离开一段时间</div>
            <div id="custom-reconnect-message">
                页面连接已断开，请刷新页面以继续使用。<br>
                您的输入内容不会丢失。
            </div>
            <button id="custom-reconnect-btn" onclick="location.reload()">
                🔄 刷新页面
            </button>
        </div>
    </div>
    
    <script>
    (function() {
        // 监控连接状态
        let lastActivity = Date.now();
        let checkInterval = null;
        let isDisconnected = false;
        
        // 更新活动时间
        function updateActivity() {
            lastActivity = Date.now();
        }
        
        // 监听用户活动
        ['click', 'keydown', 'mousemove', 'scroll', 'touchstart'].forEach(event => {
            document.addEventListener(event, updateActivity, { passive: true });
        });
        
        // 检测 Streamlit 的连接状态
        function checkConnection() {
            // 方法1: 检查是否存在 Streamlit 的连接错误元素
            const connectionError = document.querySelector(
                '[data-testid="stConnectionStatus"], ' +
                '.stException, ' +
                '[data-baseweb="modal"]'
            );
            
            // 方法2: 检查 Streamlit 的 WebSocket 状态
            const wsStatus = document.querySelector('.stStatusWidget');
            
            // 方法3: 检查是否有 "Connection error" 文本
            const pageText = document.body.innerText || '';
            const hasConnectionError = pageText.includes('Connection error') || 
                                       pageText.includes('Is Streamlit still running');
            
            // 如果检测到连接错误
            if (connectionError || hasConnectionError) {
                if (!isDisconnected) {
                    isDisconnected = true;
                    showReconnectOverlay();
                }
            }
        }
        
        // 显示重连提示
        function showReconnectOverlay() {
            const overlay = document.getElementById('custom-reconnect-overlay');
            if (overlay) {
                overlay.classList.add('show');
                
                // 隐藏原生弹窗
                hideNativeDialog();
            }
        }
        
        // 隐藏原生对话框
        function hideNativeDialog() {
            // 隐藏 Streamlit 默认的错误弹窗
            const dialogs = document.querySelectorAll(
                '[data-baseweb="modal"], ' +
                '.stException, ' +
                '[data-testid="stConnectionStatus"]'
            );
            dialogs.forEach(d => {
                d.style.display = 'none';
                d.style.visibility = 'hidden';
            });
            
            // 隐藏可能的遮罩层
            const overlays = document.querySelectorAll('[data-baseweb="modal-backdrop"]');
            overlays.forEach(o => {
                o.style.display = 'none';
            });
        }
        
        // 使用 MutationObserver 监听 DOM 变化
        const observer = new MutationObserver(function(mutations) {
            // 检查是否出现连接错误
            checkConnection();
            
            // 始终尝试隐藏原生弹窗
            if (isDisconnected) {
                hideNativeDialog();
            }
        });
        
        // 开始观察
        observer.observe(document.body, {
            childList: true,
            subtree: true,
            attributes: true
        });
        
        // 定期检查连接状态
        checkInterval = setInterval(checkConnection, 2000);
        
        // 页面可见性变化时检查
        document.addEventListener('visibilitychange', function() {
            if (document.visibilityState === 'visible') {
                // 页面重新可见时，延迟检查连接
                setTimeout(checkConnection, 1000);
            }
        });
        
        // 初始检查
        setTimeout(checkConnection, 3000);
    })();
    </script>
    """
    
    # 注入连接监控脚本
    components.html(connection_monitor_html, height=0)
    
    # 第一阶段：关键样式 + 初始隐藏（防止闪烁）
    st.markdown("""
        <style>
        /* ================================================================================== */
        /* 第一阶段：关键样式 - 立即隐藏并设置基础 */
        /* ================================================================================== */
        
        /* 立即隐藏原生导航栏 - 使用多重选择器确保覆盖 */
        [data-testid="stSidebarNav"],
        [data-testid="stSidebarNav"] *,
        nav[data-testid="stSidebarNav"] {
            display: none !important;
            visibility: hidden !important;
            height: 0 !important;
            width: 0 !important;
            overflow: hidden !important;
            position: absolute !important;
            pointer-events: none !important;
        }
        
        /* 立即设置背景色，防止白屏闪烁 */
        .stApp {
            background-color: #FFFFFF !important;
        }
        
        section[data-testid="stSidebar"] {
            background-color: #0F2B46 !important;
        }
        
        /* ================================================================================== */
        /* 全局字体与基础设置 - 使用系统字体避免加载延迟 */
        /* ================================================================================== */
        html, body, [class*="css"] {
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "Roboto", "Oxygen", "Ubuntu", "Cantarell", "Fira Sans", "Droid Sans", "Helvetica Neue", sans-serif !important;
            -webkit-font-smoothing: antialiased;
            -moz-osx-font-smoothing: grayscale;
            color: #333333;
        }
        
        /* ================================================================================== */
        /* 侧边栏样式 (DeepL 深蓝色风格) */
        /* ================================================================================== */
        section[data-testid="stSidebar"] {
            background-color: #0F2B46 !important;
            border-right: none !important;
        }
        
        /* 强制侧边栏内所有文本为白色 */
        section[data-testid="stSidebar"] .stMarkdown,
        section[data-testid="stSidebar"] h1, 
        section[data-testid="stSidebar"] h2, 
        section[data-testid="stSidebar"] h3, 
        section[data-testid="stSidebar"] p, 
        section[data-testid="stSidebar"] label,
        section[data-testid="stSidebar"] span,
        section[data-testid="stSidebar"] div {
            color: #FFFFFF !important;
        }
        
        /* 侧边栏链接/按钮 */
        section[data-testid="stSidebar"] a {
            color: #60A5FA !important;
        }
        
        /* 侧边栏分割线 */
        section[data-testid="stSidebar"] hr {
            border-color: #1E3A5F !important;
        }
        
        /* 侧边栏的信息提示框 */
        section[data-testid="stSidebar"] .stAlert {
            background-color: rgba(255, 255, 255, 0.1) !important;
            border: 1px solid rgba(255, 255, 255, 0.2) !important;
        }
        section[data-testid="stSidebar"] .stAlert * {
            color: #FFFFFF !important;
        }
        
        /* ================================================================================== */
        /* 标题颜色 */
        /* ================================================================================== */
        h1, h2, h3 {
            color: #0F2B46 !important;
            font-weight: 700 !important;
        }
        
        /* ================================================================================== */
        /* 输入框与文本域 (Card Style) */
        /* ================================================================================== */
        .stTextInput > div > div > input,
        .stTextArea > div > div > textarea {
            background-color: #F8F9FA !important;
            border: 1px solid #E5E7EB !important;
            border-radius: 8px !important;
            color: #1F2937 !important;
            box-shadow: inset 0 1px 2px rgba(0, 0, 0, 0.05) !important;
        }
        
        /* 聚焦状态 */
        .stTextInput > div > div > input:focus,
        .stTextArea > div > div > textarea:focus {
            border-color: #3B82F6 !important;
            box-shadow: 0 0 0 3px rgba(59, 130, 246, 0.2) !important;
            background-color: #FFFFFF !important;
        }
        
        /* Selectbox */
        div[data-baseweb="select"] > div {
            background-color: #F8F9FA !important;
            border-radius: 8px !important;
            border: 1px solid #E5E7EB !important;
        }

        /* ================================================================================== */
        /* 按钮样式 */
        /* ================================================================================== */
        
        /* Primary Button */
        .stButton button[kind="primary"],
        .stButton button[type="primary"],
        div[data-testid="stBaseButton-primary"],
        button[data-testid="stBaseButton-primary"],
        .stDownloadButton button,
        div[data-testid="stBaseButton-primary"] button {
            background-color: #3B82F6 !important;
            color: #FFFFFF !important;
            border: none !important;
            border-radius: 6px !important;
            font-weight: 600 !important;
            box-shadow: 0 4px 6px -1px rgba(59, 130, 246, 0.3) !important;
            transition: all 0.2s ease !important;
        }
        
        .stButton button[kind="primary"]:hover,
        .stButton button[type="primary"]:hover,
        div[data-testid="stBaseButton-primary"]:hover,
        button[data-testid="stBaseButton-primary"]:hover {
            background-color: #2563EB !important;
            box-shadow: 0 6px 8px -1px rgba(59, 130, 246, 0.4) !important;
            transform: translateY(-1px);
        }
        
        /* Secondary Button */
        .stButton button[kind="secondary"],
        .stButton button[type="secondary"],
        div[data-testid="stBaseButton-secondary"],
        button[data-testid="stBaseButton-secondary"] {
            background-color: #FFFFFF !important;
            color: #1F2937 !important;
            border: 1px solid #D1D5DB !important;
            border-radius: 6px !important;
        }

        /* ================================================================================== */
        /* 文件上传组件汉化与美化 */
        /* ================================================================================== */
        [data-testid="stFileUploader"] {
            padding: 20px !important;
            border: 2px dashed #CBD5E1 !important;
            border-radius: 12px !important;
            background-color: #F8F9FA !important;
        }
        
        /* 覆盖 "Drag and drop..." 文字 */
        [data-testid="stFileUploader"] section > div > div > span {
            visibility: hidden !important;
            position: relative !important;
        }
        [data-testid="stFileUploader"] section > div > div > span::after {
            content: "拖拽文件到此处" !important;
            visibility: visible !important;
            position: absolute !important;
            top: 0 !important;
            left: 0 !important;
            color: #0F2B46 !important;
            font-weight: 600 !important;
            font-size: 1.1em !important;
        }
        
        /* 覆盖 Limit 文字 */
        [data-testid="stFileUploader"] section > div > div > small {
            visibility: hidden !important;
            position: relative !important;
        }
        [data-testid="stFileUploader"] section > div > div > small::after {
            content: "单个文件限制 200MB • DOCX, PDF" !important;
            visibility: visible !important;
            position: absolute !important;
            top: 0 !important;
            left: 0 !important;
            color: #64748B !important;
        }
        
        /* 浏览按钮文字覆盖 */
        [data-testid="stFileUploader"] button[data-testid="stBaseButton-secondary"] {
            font-size: 0 !important;
            min-width: 100px !important;
        }
        [data-testid="stFileUploader"] button[data-testid="stBaseButton-secondary"]::after {
            content: "浏览文件" !important;
            font-size: 14px !important;
            visibility: visible !important;
            color: #1F2937 !important;
        }
        </style>
    """, unsafe_allow_html=True)


def init_page(page_title: str, page_icon: str, layout: str = "wide"):
    """
    统一的页面初始化函数
    确保配置和样式按正确顺序加载
    
    参数:
        page_title (str): 页面标题
        page_icon (str): 页面图标
        layout (str): 页面布局，默认为 "wide"
    """
    st.set_page_config(
        page_title=page_title,
        page_icon=page_icon,
        layout=layout
    )
    apply_custom_styles()


def parse_ppt_content(content: str) -> List[Dict[str, Any]]:
    """
    解析 PPT 内容文本，提取每页的标题和内容
    
    支持的格式：
    - Slide 1: 标题 / Slide 1 - 标题
    - 第一页：标题 / 第1页：标题
    - 1. 标题
    
    参数:
        content (str): PPT 内容文本
    
    返回:
        List[Dict]: 每页内容列表，格式为 [{"title": "标题", "content": ["内容1", "内容2"]}, ...]
    """
    slides = []
    
    if not content or not content.strip():
        return slides
    
    # 定义分隔模式：匹配 Slide N、第N页、数字编号等格式
    # 使用多行模式匹配
    pattern = r'(?:^|\n)(?:Slide\s*(\d+)|第([一二三四五六七八九十\d]+)页|(\d+)\s*[\.、])[\s:：\-]*(.+?)(?=(?:\nSlide\s*\d+|\n第[一二三四五六七八九十\d]+页|\n\d+\s*[\.、]|\Z))'
    
    matches = re.finditer(pattern, content, re.MULTILINE | re.DOTALL)
    
    for match in matches:
        # 提取标题部分（匹配组4）
        title_section = match.group(4).strip() if match.group(4) else ""
        
        # 分离标题和内容
        lines = title_section.split('\n')
        title = lines[0].strip() if lines else ""
        
        # 提取内容（去除标题后的所有行）
        content_lines = []
        for line in lines[1:]:
            line = line.strip()
            if line:
                # 移除列表标记（如果已有）
                if line.startswith(('-', '•', '·', '*', '1.', '2.', '3.', '4.', '5.')):
                    # 移除标记后的空格
                    line = re.sub(r'^[-•·*\d\.]\s*', '', line)
                content_lines.append(line)
        
        if title:  # 只有当标题存在时才添加幻灯片
            slides.append({
                "title": title,
                "content": content_lines
            })
    
    # 如果没有匹配到任何格式，尝试按空行分割作为备选方案
    if not slides:
        sections = re.split(r'\n\s*\n', content.strip())
        for section in sections:
            lines = [line.strip() for line in section.split('\n') if line.strip()]
            if lines:
                title = lines[0]
                content_lines = []
                for line in lines[1:]:
                    if line.startswith(('-', '•', '·', '*', '1.', '2.', '3.', '4.', '5.')):
                        line = re.sub(r'^[-•·*\d\.]\s*', '', line)
                    content_lines.append(line)
                slides.append({
                    "title": title,
                    "content": content_lines
                })
    
    return slides


def generate_pptx(slides: List[Dict[str, Any]], output_path: str = None) -> str:
    """
    根据幻灯片数据生成 PPTX 文件
    
    参数:
        slides (List[Dict]): 幻灯片数据列表
        output_path (str, optional): 输出文件路径，如果为 None 则自动生成
    
    返回:
        str: 生成的 PPTX 文件路径
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    prs = Presentation()
    
    # 设置幻灯片尺寸 (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides:
        # 使用空白布局
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.get("title", "")
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        
        # 添加内容
        content_items = slide_data.get("content", [])
        if content_items:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for i, item in enumerate(content_items):
                if i == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()
                para.text = f"• {item}" if not item.startswith(("•", "-", "·")) else item
                para.font.size = Pt(18)
                para.space_after = Pt(12)
    
    # 保存文件
    if output_path is None:
        temp_dir = Path("temp")
        temp_dir.mkdir(exist_ok=True)
        output_path = str(temp_dir / f"generated_{uuid.uuid4()}.pptx")
    
    prs.save(output_path)
    return output_path


def proofread_email(
    email_content: str,
    proofread_mode: str,
    target_language: str,
    tone: str,
    custom_terms: str = ""
) -> str:
    """
    使用 DeepSeek AI 校对邮件草稿
    
    参数:
        email_content (str): 邮件草稿内容
        proofread_mode (str): 校对模式（"仅语法修正" 或 "润色改进"）
        target_language (str): 目标语言（"中文"、"英文"、"意大利语"）
        tone (str): 语气风格（"友好"、"正式"、"简洁"）
        custom_terms (str, optional): 自定义术语/要求，默认为空字符串
    
    返回:
        str: 校对后的邮件正文
    
    异常:
        Exception: 当 API 调用失败时抛出异常
    """
    # 构建系统提示词
    # 检查是否需要使用意大利语尊称（正式语气 + 意大利语）
    italian_formal_note = ""
    if target_language == "意大利语" and tone == "正式":
        italian_formal_note = "- 必须使用意大利语正式尊称：使用 Lei（您，大写L）代替 tu，所有动词使用第三人称单数形式（如：Lei è, Lei può, Lei desidera 等），这是意大利语正式邮件的标准礼仪\n"
    
    # 构建自定义术语说明
    custom_terms_note = ""
    if custom_terms and custom_terms.strip():
        custom_terms_note = f"- 自定义术语/要求：{custom_terms.strip()}\n"
    
    system_prompt = f"""你是一位专业的邮件校对助手。
请对用户提供的邮件草稿进行校对。

校对模式：{proofread_mode}
目标语言：{target_language}
语气风格：{tone}

校对要求：
- 如果校对模式是"仅语法修正"：仅修正语法错误、拼写错误和标点错误，保持原文风格和表达方式不变
- 如果校对模式是"润色改进"：在修正语法错误的基础上，优化表达、提升流畅度和专业性，使其更符合商务邮件标准
{custom_terms_note}{italian_formal_note}- 保留邮件签名格式（如果存在）
- 不输出任何 Markdown 格式符号，包括星号（*）、下划线（_）、井号（#）等
- 仅输出校对后的邮件正文，不要添加任何说明或注释
- 保持邮件的整体结构和段落格式"""
    
    # 调用 DeepSeek API 进行校对
    return call_deepseek_api(email_content, system_prompt)
